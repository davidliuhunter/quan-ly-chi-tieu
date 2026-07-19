"""
Script tạo PowerPoint báo cáo đồ án - Quản Lý Chi Tiêu
Trường Đại học Trà Vinh
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = r"C:\Users\lieu.dong\Azure_DevOps\00.my-function\03.do-an\02.quan-ly-chi-tieu\bao-cao\ThuyetTrinh_QuanLyChiTieu_LieuSuLamDong.pptx"
IMG_DIR = r"C:\Users\lieu.dong\Azure_DevOps\00.my-function\03.do-an\02.quan-ly-chi-tieu\bao-cao\screenshots"

# ─── Colors ───────────────────────────────────────────────────────
DARK_BG = RGBColor(0x1E, 0x1B, 0x4B)      # deep indigo
ACCENT = RGBColor(0x63, 0x66, 0xF1)        # indigo accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)      # light gray
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

# ─── Helper functions ─────────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bar(slide, left=0, top=0, width=None, height=Inches(0.06)):
    if width is None: width = Inches(10)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(8)
    return txBox

def add_numbered_list(slide, left, top, width, height, items, font_size=16, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(6)
    return txBox

def add_card(slide, left, top, width, height, icon, title, value, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 icon, font_size=24, bold=False)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), Inches(0.35),
                 title, font_size=11, color=GRAY_TEXT)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.85), width - Inches(0.4), Inches(0.45),
                 value, font_size=20, color=color, bold=True)

def slide_title(slide, title, subtitle=None):
    add_accent_bar(slide, Inches(0.8), Inches(0.4), Inches(8.4))
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.7),
                 title, font_size=32, color=DARK_TEXT, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.4),
                     subtitle, font_size=14, color=GRAY_TEXT)

def add_image_safe(slide, filename, left, top, width, height=None):
    path = os.path.join(IMG_DIR, filename)
    if os.path.exists(path):
        if height:
            return slide.shapes.add_picture(path, left, top, width, height)
        else:
            return slide.shapes.add_picture(path, left, top, width=width)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height or Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        add_text_box(slide, left, top + Inches(0.5), width, Inches(0.5),
                     f"[Ảnh: {filename}]", font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
        return shape

def add_page_number(slide, num, total):
    add_text_box(slide, Inches(8.8), Inches(7.1), Inches(1), Inches(0.3),
                 f"{num}/{total}", font_size=9, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)

# ─── Create presentation ──────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
TOTAL = 18

# ===================================================================
# SLIDE 1: Title
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.6),
             "TRƯỜNG ĐẠI HỌC TRÀ VINH", font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.0), Inches(8), Inches(0.5),
             "KHOA KỸ THUẬT VÀ CÔNG NGHỆ", font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(3), Inches(2.8), Inches(4))

add_text_box(slide, Inches(0.8), Inches(3.1), Inches(8.4), Inches(1),
             "XÂY DỰNG ỨNG DỤNG WEB\nQUẢN LÝ CHI TIÊU CÁ NHÂN", 
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.5),
             "VỚI ASP.NET CORE MVC", font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.4),
             "Giảng viên hướng dẫn: ThS. Dương Ngọc Vân Khanh", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.9), Inches(8), Inches(0.4),
             "Sinh viên: Liêu Sử Lâm Đồng  |  MSSV: 170123201  |  Lớp: DX23TT9", font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.3), Inches(8), Inches(0.4),
             "Trà Vinh, tháng 7 năm 2026", font_size=13, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
add_page_number(slide, 1, TOTAL)

# ===================================================================
# SLIDE 2: Nội dung trình bày
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "NỘI DUNG TRÌNH BÀY")

agenda = [
    "1. Lý do chọn đề tài & Mục tiêu",
    "2. Công nghệ sử dụng",
    "3. Phân tích & Thiết kế hệ thống",
    "4. Cơ sở dữ liệu",
    "5. Demo các chức năng chính",
    "6. Kiểm thử & Đánh giá",
    "7. Hạn chế & Hướng phát triển",
    "8. Kết luận",
]
add_bullet_list(slide, Inches(1.2), Inches(1.7), Inches(7.5), Inches(5),
                agenda, font_size=20, color=DARK_TEXT)
add_page_number(slide, 2, TOTAL)

# ===================================================================
# SLIDE 3: Lý do chọn đề tài
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "LÝ DO CHỌN ĐỀ TÀI")

reasons = [
    "Nhu cầu quản lý tài chính cá nhân ngày càng quan trọng trong cuộc sống hiện đại",
    "Phương pháp thủ công (sổ sách, Excel) có nhiều hạn chế: dễ thất lạc, khó phân tích xu hướng",
    "Cần một ứng dụng giúp theo dõi thu chi, lập ngân sách, tiết kiệm mọi lúc mọi nơi",
    "ASP.NET Core MVC là framework mạnh mẽ, hiện đại, phù hợp để xây dựng ứng dụng web",
    "Mong muốn áp dụng kiến thức đã học vào sản phẩm thực tế có giá trị sử dụng",
]
add_bullet_list(slide, Inches(0.8), Inches(1.7), Inches(8.4), Inches(5),
                reasons, font_size=17, color=DARK_TEXT)
add_page_number(slide, 3, TOTAL)

# ===================================================================
# SLIDE 4: Mục tiêu
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "MỤC TIÊU ĐỒ ÁN")

goals = [
    "Xây dựng ứng dụng web quản lý chi tiêu cá nhân hoàn chỉnh",
    "Dashboard thống kê trực quan với biểu đồ (cột, tròn)",
    "Lập ngân sách theo danh mục, theo dõi % chi tiêu",
    "Quản lý giao dịch định kỳ & mục tiêu tiết kiệm",
    "Giao diện responsive, thân thiện, hỗ trợ Dark/Light mode",
]
add_numbered_list(slide, Inches(0.8), Inches(1.7), Inches(8.4), Inches(5),
                  goals, font_size=18, color=DARK_TEXT)

# Cards for key metrics
add_card(slide, Inches(0.8), Inches(5.6), Inches(2.6), Inches(1.3),
         "📊", "Dashboard", "Thống kê & Biểu đồ", ACCENT)
add_card(slide, Inches(3.7), Inches(5.6), Inches(2.6), Inches(1.3),
         "💳", "Giao dịch", "CRUD + Lọc & Tìm kiếm", GREEN)
add_card(slide, Inches(6.6), Inches(5.6), Inches(2.6), Inches(1.3),
         "🎯", "Ngân sách", "Hạn mức & Cảnh báo", ORANGE)
add_page_number(slide, 4, TOTAL)

# ===================================================================
# SLIDE 5: Công nghệ sử dụng
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "CÔNG NGHỆ SỬ DỤNG")

techs = [
    ("ASP.NET Core MVC 9.0", "Framework chính, kiến trúc MVC"),
    ("Entity Framework Core 9", "ORM - ánh xạ đối tượng với SQLite"),
    ("SQLite", "Database nhúng, file-based, không cần cài đặt"),
    ("Tailwind CSS", "Utility-first CSS, responsive, Dark/Light mode"),
    ("Chart.js", "Biểu đồ tương tác: Bar, Pie/Doughnut"),
    ("C# / .NET 9.0", "Ngôn ngữ lập trình chính"),
    ("Playwright", "Kiểm thử E2E tự động (17 test case)"),
]

y = Inches(1.7)
for name, desc in techs:
    add_text_box(slide, Inches(1.0), y, Inches(3.0), Inches(0.35),
                 name, font_size=16, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(4.2), y, Inches(5.0), Inches(0.35),
                 desc, font_size=14, color=GRAY_TEXT)
    y += Inches(0.55)

add_page_number(slide, 5, TOTAL)

# ===================================================================
# SLIDE 6: Kiến trúc hệ thống
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "KIẾN TRÚC HỆ THỐNG", "Mô hình MVC + Service Layer")

# Architecture diagram boxes
layers = [
    ("PRESENTATION", "Views (.cshtml)\nRazor + Tailwind CSS + Chart.js", ACCENT),
    ("CONTROLLER", "Auth · Dashboard · GiaoDich · DanhMuc\nNganSach · MucTieu · DinhKy", RGBColor(0x3B, 0x82, 0xF6)),
    ("SERVICE", "IDanhMucService · IGiaoDichService\nBusiness Logic", GREEN),
    ("DATA", "AppDbContext (EF Core)\nSQLite Database", ORANGE),
]

for i, (label, content, color) in enumerate(layers):
    y = Inches(1.7) + i * Inches(1.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y, Inches(5.5), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    
    add_text_box(slide, Inches(0.8), y + Inches(0.15), Inches(1.5), Inches(0.3),
                 label, font_size=12, color=color, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(2.8), y + Inches(0.15), Inches(5.0), Inches(0.8),
                 content, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.8), y + Inches(1.05), Inches(0.4), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY_TEXT
        arrow.line.fill.background()

add_page_number(slide, 6, TOTAL)

# ===================================================================
# SLIDE 7: ERD
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "CƠ SỞ DỮ LIỆU", "5 bảng dữ liệu · Entity Framework Core Code-First")

add_image_safe(slide, "09-erd.png", Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.5))
add_page_number(slide, 7, TOTAL)

# ===================================================================
# SLIDE 8: Mô tả bảng
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "MÔ TẢ CÁC BẢNG DỮ LIỆU")

tables_info = [
    ("📂 categories", "Danh mục thu/chi (name, type, icon, color)"),
    ("💳 transactions", "Giao dịch (amount, type, FK→category, date, description)"),
    ("🎯 budgets", "Ngân sách tháng (FK→category, month, limit_amount)"),
    ("🏆 savings_goals", "Mục tiêu tiết kiệm (target, current, deadline, progress)"),
    ("🔄 recurring", "Giao dịch định kỳ (frequency, next_date, is_active)"),
]

y = Inches(1.8)
for name, desc in tables_info:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y, Inches(8.0), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    
    add_text_box(slide, Inches(1.3), y + Inches(0.15), Inches(3.5), Inches(0.5),
                 name, font_size=16, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(4.8), y + Inches(0.15), Inches(4.0), Inches(0.5),
                 desc, font_size=14, color=GRAY_TEXT)
    y += Inches(1.0)

add_page_number(slide, 8, TOTAL)

# ===================================================================
# SLIDE 9: Use Case
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "SƠ ĐỒ USE CASE")

add_image_safe(slide, "10-usecase.png", Inches(1.0), Inches(1.6), Inches(8.0), Inches(5.5))
add_page_number(slide, 9, TOTAL)

# ===================================================================
# SLIDE 10: Demo - Đăng nhập
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "DEMO: ĐĂNG NHẬP", "Session-based Authentication · Cookie HttpOnly")

add_image_safe(slide, "00-login.png", Inches(1.5), Inches(1.6), Inches(7.0), Inches(4.0))

add_text_box(slide, Inches(0.8), Inches(5.9), Inches(8.4), Inches(0.4),
             "🔐 Mật khẩu hardcoded → Session \"IsAdmin\" → Bảo vệ tất cả route /quan-tri/*",
             font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(8.4), Inches(0.4),
             "⏱️ Idle timeout 4 giờ · Cookie HttpOnly chống XSS",
             font_size=13, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 10, TOTAL)

# ===================================================================
# SLIDE 11: Demo - Dashboard
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "DEMO: DASHBOARD TỔNG QUAN", "Thẻ thống kê · Biểu đồ Chart.js · Animated Counter")

add_image_safe(slide, "01-dashboard.png", Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.0))

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.4),
             "📊 Biểu đồ cột: Thu nhập & Chi tiêu 6 tháng | 🍩 Biểu đồ tròn: Phân bổ theo danh mục",
             font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.4),
             "💚 Animated Counter với requestAnimationFrame · Giao dịch gần đây (top 5)",
             font_size=13, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 11, TOTAL)

# ===================================================================
# SLIDE 12: Demo - Giao dịch
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "DEMO: QUẢN LÝ GIAO DỊCH", "CRUD đầy đủ · Bộ lọc thông minh · Form validation")

add_image_safe(slide, "02-giao-dich.png", Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.0))

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.4),
             "🔍 Lọc: loại (thu/chi) · từ khóa · khoảng ngày · khoảng tiền",
             font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.4),
             "✅ Validation client + server (Data Annotations) · Toast thông báo",
             font_size=13, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 12, TOTAL)

# ===================================================================
# SLIDE 13: Demo - Ngân sách & Mục tiêu
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "DEMO: NGÂN SÁCH & MỤC TIÊU TIẾT KIỆM")

add_image_safe(slide, "03-ngan-sach.png", Inches(0.3), Inches(1.5), Inches(4.6), Inches(2.8))
add_image_safe(slide, "04-muc-tieu.png", Inches(5.2), Inches(1.5), Inches(4.6), Inches(2.8))

add_text_box(slide, Inches(0.3), Inches(4.5), Inches(4.6), Inches(0.4),
             "🎯 Ngân sách: hạn mức · % chi tiêu · cảnh báo màu",
             font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.2), Inches(4.5), Inches(4.6), Inches(0.4),
             "🏆 Mục tiêu: tiến độ · deadline · tự động hoàn thành",
             font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.3), Inches(4.8), Inches(4.6), Inches(0.3),
             "🟢 <80% · 🟡 80-100% · 🔴 >100% vượt hạn mức",
             font_size=11, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.2), Inches(4.8), Inches(4.6), Inches(0.3),
             "Tự động đánh dấu is_completed khi đạt target",
             font_size=11, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Bottom: Định kỳ & Danh mục
add_image_safe(slide, "05-dinh-ky.png", Inches(0.3), Inches(5.1), Inches(4.6), Inches(2.0))
add_image_safe(slide, "06-danh-muc.png", Inches(5.2), Inches(5.1), Inches(4.6), Inches(2.0))
add_page_number(slide, 13, TOTAL)

# ===================================================================
# SLIDE 14: Demo - Mobile
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "DEMO: GIAO DIỆN MOBILE", "Responsive Design · Bottom Navigation · PWA Ready")

add_image_safe(slide, "07-mobile-dashboard.png", Inches(2.5), Inches(1.5), Inches(5.0), Inches(5.5))
add_page_number(slide, 14, TOTAL)

# ===================================================================
# SLIDE 15: Flowchart
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "LƯU ĐỒ XỬ LÝ", "Đăng nhập & Thêm giao dịch")

add_image_safe(slide, "11-login-flow.png", Inches(0.3), Inches(1.5), Inches(4.6), Inches(5.5))
add_image_safe(slide, "12-transaction-flow.png", Inches(5.2), Inches(1.5), Inches(4.6), Inches(5.5))
add_page_number(slide, 15, TOTAL)

# ===================================================================
# SLIDE 16: Kiểm thử
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "KIỂM THỬ HỆ THỐNG", "17 test case · Playwright E2E · Tỷ lệ Pass: 100%")

test_cases = [
    "✅ TC01-03: Đăng nhập (sai mật khẩu, đúng mật khẩu, redirect)",
    "✅ TC04-05: Dashboard (thẻ thống kê, biểu đồ Chart.js)",
    "✅ TC06-10: CRUD Giao dịch (thêm, sửa, xóa, lọc, tìm kiếm)",
    "✅ TC11: CRUD Danh mục (thêm, sửa, xóa)",
    "✅ TC12-13: Ngân sách & Mục tiêu tiết kiệm",
    "✅ TC14: Giao dịch định kỳ (4 loại frequency)",
    "✅ TC15-17: Dark/Light mode · Responsive · Đăng xuất",
]

add_bullet_list(slide, Inches(0.8), Inches(1.7), Inches(8.4), Inches(4.5),
                test_cases, font_size=16, color=DARK_TEXT)

# Result badge
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(6.3), Inches(4.0), Inches(0.6))
shape.fill.solid()
shape.fill.fore_color.rgb = GREEN
shape.line.fill.background()
add_text_box(slide, Inches(3.0), Inches(6.35), Inches(4.0), Inches(0.5),
             "✅ 17/17 TEST CASE PASS (100%)", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 16, TOTAL)

# ===================================================================
# SLIDE 17: Kết quả đạt được
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "KẾT QUẢ ĐẠT ĐƯỢC")

results = [
    "Xây dựng thành công ứng dụng web với 15 chức năng chính",
    "Dashboard thống kê với 2 loại biểu đồ (Bar + Pie/Doughnut)",
    "Kiến trúc MVC chuẩn + Service Layer + Dependency Injection",
    "EF Core Code-First với 5 bảng dữ liệu SQLite",
    "Session-based Authentication bảo vệ toàn bộ trang quản lý",
    "17 test case Playwright E2E - Pass 100%",
    "Responsive Design: Desktop (sidebar) + Mobile (bottom nav)",
    "Dark/Light mode · Animated counter · Toast notification",
]

add_numbered_list(slide, Inches(0.8), Inches(1.7), Inches(8.4), Inches(5.0),
                  results, font_size=16, color=DARK_TEXT)
add_page_number(slide, 17, TOTAL)

# ===================================================================
# SLIDE 18: Hạn chế & Hướng phát triển
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_title(slide, "HẠN CHẾ & HƯỚNG PHÁT TRIỂN")

# Left: Hạn chế
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(4.0), Inches(0.4),
             "⚠️ HẠN CHẾ", font_size=20, color=RED, bold=True)

limits = [
    "Xác thực đơn giản (single user, mật khẩu hardcoded)",
    "SQLite - khó mở rộng khi nhiều người dùng",
    "Chưa có xuất báo cáo PDF/Excel",
    "Chưa có đồng bộ ngân hàng hay app mobile",
    "Tailwind CSS qua CDN (chưa build tối ưu)",
    "Chưa hỗ trợ đa ngôn ngữ (i18n)",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(4.2), Inches(4.5),
                limits, font_size=13, color=DARK_TEXT)

# Right: Hướng phát triển
add_text_box(slide, Inches(5.5), Inches(1.6), Inches(4.0), Inches(0.4),
             "🚀 HƯỚNG PHÁT TRIỂN", font_size=20, color=GREEN, bold=True)

future = [
    "ASP.NET Core Identity: đa người dùng, phân quyền",
    "Nâng cấp SQL Server / PostgreSQL cho scale",
    "Xuất báo cáo PDF/Excel, Web API cho mobile app",
    "Tích hợp AI gợi ý tiết kiệm, dự đoán xu hướng",
    "Build Tailwind CSS qua npm, tối ưu CSS",
    "Triển khai Azure / Vercel production",
]
add_bullet_list(slide, Inches(5.5), Inches(2.1), Inches(4.2), Inches(4.5),
                future, font_size=13, color=DARK_TEXT)
add_page_number(slide, 18, TOTAL)

# ===================================================================
# SLIDE 19: Kết luận
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.8),
             "KẾT LUẬN", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_bar(slide, Inches(3), Inches(2.3), Inches(4))

conclusion = [
    "Đồ án đã hoàn thành đầy đủ các mục tiêu đề ra",
    "Ứng dụng web quản lý chi tiêu cá nhân với 15 chức năng",
    "Kiến trúc MVC chuẩn, sẵn sàng mở rộng",
    "17/17 test case Playwright E2E pass 100%",
    "Nền tảng vững chắc cho các phát triển tiếp theo",
]

add_bullet_list(slide, Inches(1.5), Inches(2.8), Inches(7), Inches(3.5),
                conclusion, font_size=20, color=WHITE)

add_page_number(slide, 19, TOTAL)

# ===================================================================
# SLIDE 20: Cảm ơn
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(1), Inches(2.2), Inches(8), Inches(1.2),
             "CẢM ƠN THẦY CÔ\nĐÃ LẮNG NGHE", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_bar(slide, Inches(3), Inches(3.8), Inches(4))

add_text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.6),
             "XÂY DỰNG ỨNG DỤNG WEB QUẢN LÝ CHI TIÊU CÁ NHÂN", 
             font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.4),
             "Sinh viên: Liêu Sử Lâm Đồng | MSSV: 170123201 | Lớp: DX23TT9",
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.9), Inches(8), Inches(0.4),
             "GVHD: ThS. Dương Ngọc Vân Khanh",
             font_size=15, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.4), Inches(8), Inches(0.4),
             "Trà Vinh, tháng 7 năm 2026",
             font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)

add_page_number(slide, 20, TOTAL)

# ─── Save ─────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"✅ PowerPoint đã tạo: {OUTPUT}")
print(f"   {len(prs.slides)} slides · {os.path.getsize(OUTPUT)} bytes")
